#!/usr/bin/env python3
import os
from pptx import Presentation
from reco.models import DeckData

def build_ppt(deck: DeckData, output_path: str):
    template_path = "pptx_generator/templates/base.pptx"
    prs = Presentation(template_path) if os.path.exists(template_path) else Presentation()

    # 1. Brief Reminder
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "1. Brief Reminder"
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for obj in deck.brief_reminder.objectives:
        p = tf.add_paragraph()
        p.text = f"• {obj}"
    p = tf.add_paragraph()
    p.text = deck.brief_reminder.internal_reformulation

    # 2. Brand Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "2. Brand Overview"
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for paragraph in deck.brand_overview.description_paragraphs:
        p = tf.add_paragraph()
        p.text = paragraph

    # 3. State of Play
    for idx, section in enumerate(deck.state_of_play, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"3.{idx} State of Play – {section.theme}"
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        for ev in section.evidence:
            p = tf.add_paragraph()
            p.text = f"• {ev}"

    # 4. Ideas
    for idx, idea in enumerate(deck.ideas, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"4. Idea #{idx}: {idea.label}"
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        for bullet in idea.bullets:
            p = tf.add_paragraph()
            p.text = f"• {bullet}"

    # 5. Timeline
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "5. Timeline"
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for m in deck.timeline:
        p = tf.add_paragraph()
        p.text = f"{m.deadline}: {m.label}"

    # 6. Budget
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "6. Budget"
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for b in deck.budget:
        p = tf.add_paragraph()
        p.text = f"{b.category}: €{b.estimate} ({b.comment})"

    prs.save(output_path)
