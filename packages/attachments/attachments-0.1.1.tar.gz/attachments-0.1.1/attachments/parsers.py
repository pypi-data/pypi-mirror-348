"""File parsing logic."""

import os
from abc import ABC, abstractmethod
from .exceptions import ParsingError
from .utils import parse_index_string, parse_image_operations, parse_audio_operations # Added parse_audio_operations
import io # Ensure io is imported for BytesIO if AudioParser exports to it directly

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import pillow_heif
    pillow_heif.register_heif_opener()
except ImportError:
    print("Warning: pillow-heif not installed. HEIC/HEIF support will be unavailable.")
    pass # Allow the rest of the module to load

try:
    from pydub import AudioSegment
    from pydub.exceptions import CouldntDecodeError
except ImportError:
    AudioSegment = None # pydub is optional, handled in AudioParser
    CouldntDecodeError = None
    print("Warning: pydub not installed. Audio processing capabilities will be limited. pip install pydub")

class BaseParser(ABC):
    """Abstract base class for file parsers."""
    @abstractmethod
    def parse(self, file_path, indices=None):
        """Parses the file and returns structured content."""
        pass

class ParserRegistry:
    """Manages registration and retrieval of parsers."""
    def __init__(self):
        self.parsers = {}

    def register(self, type_name, parser_instance):
        """Registers a parser instance for a given type name."""
        if not isinstance(parser_instance, BaseParser):
            raise TypeError("Parser instance must be a subclass of BaseParser.")
        self.parsers[type_name] = parser_instance

    def get_parser(self, type_name):
        """Retrieves a registered parser by type name."""
        parser = self.parsers.get(type_name)
        if not parser:
            raise ValueError(f"No parser registered for type '{type_name}'.")
        return parser

class PDFParser(BaseParser):
    """Parses PDF files using PyMuPDF (fitz)."""
    def parse(self, file_path, indices=None):
        """Parses the PDF file and extracts text and image information.
        `indices` can be a string (e.g., "1,3-5,N") specifying page numbers or ranges.
        """
        if fitz is None:
            raise ParsingError("PyMuPDF (fitz) is not installed. Please install it to parse PDF files. You can typically install it with: pip install PyMuPDF")

        try:
            doc = fitz.open(file_path)
            text_parts = []
            images_info = []
            
            num_pages_total = doc.page_count
            pages_to_process_indices = [] # 0-indexed list

            if indices and isinstance(indices, str):
                pages_to_process_indices = parse_index_string(indices, num_pages_total)
                if not pages_to_process_indices and indices.strip(): # If string was not empty but parsing yielded no pages
                    # This could be due to invalid indices like "999" for a 10-page doc
                    # Or a malformed string like "abc"
                    # parse_index_string prints warnings for unparseable parts.
                    # If the result is empty, it means no valid pages were selected.
                    # We could choose to process all, or none, or raise error.
                    # For now, let's process no pages if specific indices were given but resulted in empty set.
                    print(f"Warning: PDF index string '{indices}' resulted in no pages to process for {file_path}. No content will be extracted.")
            elif isinstance(indices, list): # Support direct list of 0-indexed integers (internal use?)
                pages_to_process_indices = [p for p in indices if 0 <= p < num_pages_total]
            else: # No indices or unhandled type, process all pages
                pages_to_process_indices = list(range(num_pages_total))
            
            if not pages_to_process_indices and num_pages_total > 0 and indices:
                # If indices were provided, but resulted in an empty list, extract nothing.
                # (This condition is now partly handled above, but good for clarity)
                pass # text_parts and images_info will remain empty
            elif not pages_to_process_indices and num_pages_total > 0 and not indices:
                 # If no indices provided and pages exist, process all (already default by range)
                 pages_to_process_indices = list(range(num_pages_total))
            elif not pages_to_process_indices and num_pages_total == 0:
                # No pages in doc, nothing to process
                pass 

            for page_num_0_indexed in pages_to_process_indices:
                # Page numbers are already validated by parse_index_string to be within bounds
                page = doc.load_page(page_num_0_indexed)
                text_parts.append(page.get_text("text").strip())
                
                img_list = page.get_images(full=True)
                for img_index, img in enumerate(img_list):
                    xref = img[0]
                    images_info.append({
                        "page_num": page_num_0_indexed + 1, # Store as 1-indexed for consistency in output
                        "img_index_on_page": img_index,
                        "xref": xref,
                        "width": img[2],
                        "height": img[3],
                        "format": img[5],
                        "bbox": page.get_image_bbox(img).irect.to_json()
                    })
            
            doc.close()
            text_content = "\n\n".join(text_parts)

            return {
                "text": text_content,
                "images": images_info,
                "num_pages": num_pages_total, # Total pages in doc
                "file_path": file_path,
                # Store 1-indexed pages that were actually processed
                "indices_processed": [p + 1 for p in pages_to_process_indices]
            }
        except FileNotFoundError:
            raise ParsingError(f"Error parsing PDF: File not found at {file_path}")
        except Exception as e:
            if "no such file or directory" in str(e).lower() or \
               "cannot open" in str(e).lower():
                raise ParsingError(f"Error parsing PDF: File not found or cannot be opened at {file_path}. (PyMuPDF: {e})")
            elif "damaged" in str(e).lower() or "cannot be opened" in str(e).lower():
                 raise ParsingError(f"Error reading PDF file {file_path}. The file might be corrupted or encrypted. (PyMuPDF: {e})")
            raise ParsingError(f"An unexpected error occurred while parsing PDF file {file_path} with PyMuPDF: {e}")

class PPTXParser(BaseParser):
    """Parses PowerPoint (PPTX) files."""
    def parse(self, file_path, indices=None):
        """Parses the PPTX file and extracts text content from selected slides.
        `indices` can be a string (e.g., "1,3-5,N") specifying slide numbers or ranges.
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            all_slides = list(prs.slides) # Convert to list to use indices
            num_slides_total = len(all_slides)
            slides_to_process_indices = [] # 0-indexed list

            if indices and isinstance(indices, str):
                slides_to_process_indices = parse_index_string(indices, num_slides_total)
                if not slides_to_process_indices and indices.strip():
                    print(f"Warning: PPTX index string '{indices}' resulted in no slides to process for {file_path}. No content will be extracted.")
            elif isinstance(indices, list):
                slides_to_process_indices = [s for s in indices if 0 <= s < num_slides_total]
            else: # No indices or unhandled type, process all slides
                slides_to_process_indices = list(range(num_slides_total))

            text_content_parts = []
            processed_slide_numbers_1_indexed = []

            if not slides_to_process_indices and num_slides_total > 0 and indices:
                 pass # Extract nothing if indices given but empty result
            elif not slides_to_process_indices and num_slides_total > 0 and not indices:
                 slides_to_process_indices = list(range(num_slides_total)) # Process all if no indices
            elif not slides_to_process_indices and num_slides_total == 0:
                 pass # No slides in doc

            for slide_num_0_indexed in slides_to_process_indices:
                slide = all_slides[slide_num_0_indexed]
                slide_text_parts = []
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            slide_text_parts.append(run.text)
                
                if slide_text_parts:
                    # Using original 1-indexed number for slide header for clarity
                    text_content_parts.append(f"--- Slide {slide_num_0_indexed + 1} ---\n{' '.join(slide_text_parts)}")
                else: # Add a marker for blank slides if they were selected
                    text_content_parts.append(f"--- Slide {slide_num_0_indexed + 1} ---\n[Blank Slide or No Text Content]")
                processed_slide_numbers_1_indexed.append(slide_num_0_indexed + 1)

            final_text_content = "\n\n".join(text_content_parts).strip()

            return {
                "text": final_text_content,
                "num_slides": num_slides_total, # Total slides in doc
                "file_path": file_path,
                "indices_processed": processed_slide_numbers_1_indexed # 1-indexed slides processed
            }
        except ImportError:
            raise ParsingError("python-pptx is not installed. Please install it to parse PPTX files. You can typically install it with: pip install python-pptx")
        except FileNotFoundError:
            raise ParsingError(f"Error parsing PPTX: File not found at {file_path}")
        except Exception as e:
            # Check if it's a PackageNotFoundError, which might indicate a corrupted file for python-pptx
            if "PackageNotFoundError" in str(type(e)):
                 raise ParsingError(f"Error parsing PPTX file {file_path}. The file might be corrupted, not a valid PPTX, or an issue with python-pptx: {e}")
            raise ParsingError(f"An unexpected error occurred while parsing PPTX file {file_path}: {e}")

class HTMLParser(BaseParser):
    """Parses HTML files into Markdown text using html2text."""
    def parse(self, file_path, indices=None):
        """Parses the HTML file and converts its content to Markdown.
        `indices` is currently ignored for HTML files.
        """
        try:
            import html2text
            h = html2text.HTML2Text()
            # Configure html2text options if needed, e.g.:
            # h.ignore_links = True
            # h.ignore_images = True
            # h.body_width = 0 # No line wrapping

            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            markdown_content = h.handle(html_content)
            
            # For HTML, num_pages/num_slides isn't directly applicable in the same way.
            # We can omit it or add other relevant metadata if available (e.g., title).
            # For now, keep it simple.
            return {
                "text": markdown_content.strip(),
                "file_path": file_path,
                # "indices_processed": [] # Or None, as indices are not used yet
            }
        except ImportError:
            raise ParsingError("html2text is not installed. Please install it to parse HTML files. You can typically install it with: pip install html2text")
        except FileNotFoundError:
            raise ParsingError(f"Error parsing HTML: File not found at {file_path}")
        except Exception as e:
            raise ParsingError(f"An unexpected error occurred while parsing HTML file {file_path}: {e}")

class ImageParser(BaseParser):
    """Parses image files using Pillow, supporting transformations."""
    def parse(self, file_path, indices=None):
        """Parses the image file, extracts metadata, applies transformations, 
        and stores the Pillow Image object.
        `indices` for images is an operation string, e.g., "resize:100x100,rotate:90"
        """
        try:
            from PIL import Image, UnidentifiedImageError, ImageOps
        except ImportError:
            raise ParsingError("Pillow (PIL) is not installed. Please install it to parse image files. You can typically install it with: pip install Pillow")

        operations = parse_image_operations(indices) # `indices` is the ops_str for images
        
        try:
            img = Image.open(file_path)
            img.load() 
            original_format_val = img.format # Capture before any img reassignment
            original_mode_val = img.mode   # Capture before any img reassignment

            # Apply operations
            if 'rotate' in operations:
                angle = operations['rotate']
                # Pillow's rotate can expand the image. For 90, 180, 270, we want exact rotations.
                # For 0, no change. For 180, it's fine.
                # For 90 and 270, using transpose is better to avoid black borders/resizing.
                if angle == 90:
                    img = img.transpose(Image.Transpose.ROTATE_90)
                elif angle == 180:
                    img = img.rotate(180) # rotate(180) is exact
                elif angle == 270:
                    img = img.transpose(Image.Transpose.ROTATE_270)
                # Not applying img.rotate(angle, expand=True) for other angles for now to keep it simple
                # and avoid unexpected dimension changes unless explicitly handled.

            if 'resize' in operations:
                target_w, target_h = operations['resize']
                original_w, original_h = img.width, img.height
                
                if target_w is None and target_h is None: # Should be caught by parser, but defensive
                    pass # No resize
                elif target_w is None: # auto width, fixed height
                    aspect_ratio = original_w / original_h
                    target_w = int(target_h * aspect_ratio)
                elif target_h is None: # auto height, fixed width
                    aspect_ratio = original_h / original_w
                    target_h = int(target_w * aspect_ratio)
                
                # Now target_w and target_h are determined
                if (target_w, target_h) != (original_w, original_h) and target_w > 0 and target_h > 0:
                    img = img.resize((target_w, target_h), Image.Resampling.LANCZOS)
            
            # Original mode conversion logic (P, PA, L)
            if img.mode in ('P', 'PA'):
                if 'A' in img.mode or (img.info.get('transparency') is not None):
                    img = img.convert('RGBA')
                else:
                    img = img.convert('RGB')
            elif img.mode == 'L':
                img = img.convert('RGB')

            # Determine output format and quality (can be overridden by operations)
            output_format = operations.get('format')
            if not output_format: # If not specified by operation, infer a sensible default
                if img.format and img.format.upper() in ['PNG', 'GIF', 'WEBP'] and img.mode == 'RGBA':
                    output_format = 'png' # Preserve alpha for these types if present
                else:
                    output_format = 'jpeg' # Default to JPEG for others or if no alpha
            
            output_quality = operations.get('quality', 90) # Default quality

            # Use the captured original_format_val and original_mode_val for the text representation
            text_representation = f"[Image: {os.path.basename(file_path)} (original: {original_format_val} {original_mode_val}) -> processed to {img.width}x{img.height} for output as {output_format}]"
            if indices: # `indices` here is the ops_str
                 text_representation = f"[Image: {os.path.basename(file_path)} (original: {original_format_val} {original_mode_val}, ops: \"{indices}\") -> processed to {img.width}x{img.height} for output as {output_format}]"

            return {
                "text": text_representation,
                "file_path": file_path,
                "image_object": img,  # This is now the (potentially) transformed image
                "width": img.width,   # Width of the transformed image
                "height": img.height, # Height of the transformed image
                "original_format": original_format_val, # Use captured value
                "original_mode": original_mode_val,     # Use captured value
                "output_format": output_format, # Target format for base64/saving
                "output_quality": output_quality, # Target quality for base64/saving
                "applied_operations": operations # Store the operations that were parsed
            }
        except FileNotFoundError:
            raise ParsingError(f"Error parsing image: File not found at {file_path}")
        except UnidentifiedImageError:
            raise ParsingError(f"Error parsing image: Cannot identify image file at {file_path}. It might be corrupted or not a supported image format by Pillow.")
        except Exception as e:
            raise ParsingError(f"An unexpected error occurred while parsing image file {file_path} with Pillow: {e}")

# Audio Parser
class AudioParser(BaseParser):
    """Parses and processes audio files using pydub for API submission preparation."""
    def parse(self, file_path, indices=None):
        """Loads audio, applies transformations (format, samplerate, channels, bitrate),
        and returns a dictionary with the processed audio segment and metadata.
        Default processing: convert to 16kHz mono WAV if no operations specified.
        """
        if AudioSegment is None:
            raise ParsingError("pydub is not installed. Please install it to enable audio processing: pip install pydub")

        if not os.path.exists(file_path):
            raise ParsingError(f"Audio file not found: {file_path}")

        operations = parse_audio_operations(indices if isinstance(indices, str) else "")
        
        # Default operations if not specified by user
        # Whisper prefers mono 16kHz WAV for best performance.
        # Let's make this our default if no format is specified by user.
        output_format = operations.get('format', 'wav').lower()
        target_samplerate = operations.get('samplerate') 
        target_channels = operations.get('channels')
        target_bitrate = operations.get('bitrate')

        if not operations.get('format') and not target_samplerate and not target_channels: # If no user ops affecting these
            target_samplerate = target_samplerate or 16000
            target_channels = target_channels or 1 # Mono

        original_basename = os.path.basename(file_path)
        original_name_part, _ = os.path.splitext(original_basename)
        
        try:
            # Load the audio file. pydub usually infers input format automatically.
            # Explicitly providing format might be useful if detection fails.
            # For now, rely on pydub's auto-detection.
            try:
                audio_segment = AudioSegment.from_file(file_path)
            except CouldntDecodeError as e:
                # Attempt with common extensions if pydub fails to guess format
                # (This is a fallback, usually pydub is good with extensions)
                _, ext = os.path.splitext(file_path)
                ext = ext.lower().replace('.', '')
                if ext in ['mp3', 'wav', 'flac', 'ogg', 'm4a', 'aac', 'opus']:
                    try:
                        print(f"Retrying pydub load for {file_path} with explicit format: {ext}")
                        audio_segment = AudioSegment.from_file(file_path, format=ext)
                    except CouldntDecodeError:
                         raise ParsingError(f"pydub could not decode audio file {file_path} even with explicit format '{ext}': {e}") from e
                else:
                    raise ParsingError(f"pydub could not decode audio file {file_path}: {e}") from e
            except Exception as e_load: # Catch other pydub loading errors
                raise ParsingError(f"Error loading audio file {file_path} with pydub: {e_load}") from e_load


            # Apply transformations
            if target_samplerate and audio_segment.frame_rate != target_samplerate:
                audio_segment = audio_segment.set_frame_rate(target_samplerate)
            
            if target_channels and audio_segment.channels != target_channels:
                if target_channels == 1:
                    audio_segment = audio_segment.set_channels(1) # Convert to mono
                # else: pydub doesn't have a simple "set_channels(2)" if it's already stereo.
                # If converting mono to stereo, it's more complex (e.g. duplicate channel).
                # For STT, mono is usually preferred or sufficient.
                # If user explicitly asks for stereo and it's already stereo, no change.
                # If user explicitly asks for stereo and it's mono, this is an edge case to consider.
                # For now, primary use case is ensuring mono if target_channels is 1.

            # Note: format conversion and bitrate are handled during export by pydub.
            # We store the target format and bitrate for the .audios property to use.

            processed_filename_for_api = f"{original_name_part}.{output_format}"
            
            text_representation_parts = [f"Audio: {original_basename}"]
            if indices:
                text_representation_parts.append(f"ops: \"{indices}\"")
            
            text_representation_parts.append(f"-> processed to {output_format}")
            if target_samplerate: text_representation_parts.append(f"{target_samplerate//1000}kHz")
            if target_channels == 1: text_representation_parts.append("mono")
            elif target_channels == 2: text_representation_parts.append("stereo")
            
            text_representation = f"[{' '.join(text_representation_parts)}]"
            
            return {
                "text": text_representation,
                "raw_path": file_path, # Keep original path for reference
                "audio_segment": audio_segment, # The pydub AudioSegment object
                "original_basename": original_basename, # For reference
                "processed_filename_for_api": processed_filename_for_api, # e.g., "input.wav"
                "output_format": output_format, # e.g., "wav"
                "output_samplerate": audio_segment.frame_rate, # Actual rate after processing
                "output_channels": audio_segment.channels,   # Actual channels after processing
                "output_bitrate": target_bitrate, # Requested bitrate for export (string like "128k")
                "applied_operations": operations # Store parsed operations
            }

        except CouldntDecodeError as e: # Should be caught above, but as a safety net
            raise ParsingError(f"pydub failed to decode audio file: {file_path}. Ensure ffmpeg/libav is installed if it's not a WAV/MP3. Error: {e}")
        except FileNotFoundError: # Should be caught by os.path.exists already
            raise ParsingError(f"Audio file not found during pydub processing (should have been caught earlier): {file_path}")
        except Exception as e:
            raise ParsingError(f"An unexpected error occurred while processing audio file {file_path} with pydub: {e}")

# Example of how parsers might be registered (this would typically happen in the Attachments core or user code)
# parser_registry = ParserRegistry()
# parser_registry.register('pdf', PDFParser())
# parser_registry.register('pptx', PPTXParser()) 