from pptx import Presentation
from pptx.util import Inches
import os

# Create a new presentation
prs = Presentation()

# Slide 1: Title Slide
slide_layout_1 = prs.slide_layouts[0]  # Title Slide Layout
slide_1 = prs.slides.add_slide(slide_layout_1)
title_1 = slide_1.shapes.title
title_1.text = "Slide 1 Title"
subtitle_1 = slide_1.placeholders[1]
subtitle_1.text = "This is the first slide.\nContent for page 1."

# Slide 2: Title and Content
slide_layout_2 = prs.slide_layouts[1]  # Title and Content Layout
slide_2 = prs.slides.add_slide(slide_layout_2)
title_2 = slide_2.shapes.title
title_2.text = "Slide 2 Title"
body_2 = slide_2.placeholders[1]
tf_2 = body_2.text_frame
tf_2.text = "This is the second slide."
p_2 = tf_2.add_paragraph()
p_2.text = "Content for page 2."

# Slide 3: Blank slide with a textbox
slide_layout_3 = prs.slide_layouts[5]  # Blank Layout
slide_3 = prs.slides.add_slide(slide_layout_3)
left = top = width = height = Inches(1.0)
txBox = slide_3.shapes.add_textbox(left, top, width*5, height)
tf_3 = txBox.text_frame
tf_3.text = "Slide 3 Title\nThis is the third slide.\nContent for page 3."

# Save the presentation
output_filename = "sample.pptx"
# Ensure the script saves it in its own directory (tests/test_data/)
script_dir = os.path.dirname(__file__)
output_path = os.path.join(script_dir, output_filename)

prs.save(output_path)
print(f"Successfully created {output_path} with 3 slides.") 