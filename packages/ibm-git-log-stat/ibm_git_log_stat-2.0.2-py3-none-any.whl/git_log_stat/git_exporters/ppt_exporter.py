import textwrap
from collections import defaultdict
from logging import Logger

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Inches

from git_log_stat.git_exporters.base_exporter import BaseExporter


def generate_commit_summary_pptx(log: Logger, commit_logs: str | list[str],
                                 pr_logs: str | list[str], output_file="commit_summary.pptx",
                                 skip_summary=False):
    """
    Generate PPT file with commits, pr if present and summary if requested.
    :param log: logger to use for logging during pdf creation
    :param commit_logs: commit logs to be written and processed through NLP
    :param pr_logs: pull request logs to be written to pdf
    :param output_pdf_path: output pdf file path
    :param skip_summary: boolean flag to exclude summary. Default is False.
    :return: Path where pdf file is generated
    """
    prs = Presentation()

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Commit Summary Report"
    slide.placeholders[1].text = "Detailed Report of Commits and Pull Requests"

    # Natural Language Summary
    if not skip_summary:
        try:
            from git_log_stat.git_repo.git_nlp import init_summarizer
            summarizer = init_summarizer()
        except ImportError:
            summarizer = None

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Natural Language Summary"

        log.info("🔄 Summarizing commit logs...")

        combined_text = ". ".join([line.split("|", 2)[-1].strip() for line in commit_logs.split("\n")])
        combined_text = " ".join(combined_text.split()[:1024])
        summary = summarizer(combined_text, max_length=255, min_length=40, do_sample=False)[0]['summary_text']

        text_box = slide.shapes.placeholders[1]
        text_frame = text_box.text_frame
        for line in textwrap.wrap(summary, width=100):
            p = text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(12)

    # Commits by Author
    grouped_commits = defaultdict(list)
    for line in commit_logs.split("\n"):
        parts = [p.strip() for p in line.split('|')]
        if len(parts) >= 4:
            date, author, commit_id, *message_parts = parts
            message = " | ".join(message_parts)
            grouped_commits[author].append((date, commit_id, message))

    for author, commits in grouped_commits.items():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Commits by {author}"
        text_frame = slide.shapes.placeholders[1].text_frame
        text_frame.clear()

        for date, commit_id, message in commits:
            full_text = f"{date} | {commit_id} | {message}"
            wrapped_lines = textwrap.wrap(full_text, width=100)
            for line in wrapped_lines:
                p = text_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(10)
                p.alignment = PP_ALIGN.LEFT

    # PR Logs Slide
    if pr_logs:
        if isinstance(pr_logs, str):
            pr_logs = pr_logs.strip().splitlines()

        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        slide.shapes.title.text = "Pull Request Summary"

        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)

        shape = slide.shapes.add_textbox(left, top, width, height)
        text_frame = shape.text_frame

        for pr in pr_logs:
            for line in textwrap.wrap(pr.strip(), width=100):
                p = text_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(11)

    # Save file
    prs.save(output_file)
    log.info(f"✅ PowerPoint file saved to {output_file}")


class PptExporter(BaseExporter):
    def export(self, output_file_name, commit_output: str | list[str], pr_output: str | list[str] = None,
               skip_summary=False):
        self.log.info("Generating PPT %s", output_file_name)
        generate_commit_summary_pptx(
            self.log,
            commit_logs=commit_output,
            pr_logs=pr_output,
            output_file=output_file_name,
            skip_summary=skip_summary
        )
