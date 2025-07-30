import logging
from logging import Logger
from typing import List, Dict, Any

import requests
from docx import Document
from openpyxl.utils import get_column_letter
from openpyxl.workbook import Workbook
from openpyxl.worksheet.worksheet import Worksheet
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Table, TableStyle, Paragraph, Spacer, SimpleDocTemplate


class JiraExporter:
    def __init__(self, log: Logger, jira_url: str, auth: tuple):
        self.log = log
        self.jira_url = jira_url.rstrip('/')
        self.auth = auth

    def generate_jira_report(self, user_email: str, jql: str = None, start_date: str = None, end_date: str = None) -> List[Dict[str, Any]]:
        """Returns list of dicts with issue_id, status, assignee, reporter, summary"""
        try:
            date_filters = []
            if start_date:
                date_filters.append(f'updated >= "{start_date}"')
            if end_date:
                date_filters.append(f'updated <= "{end_date}"')

            # Build final JQL
            default_jql = f'assignee = "{user_email}"'
            if date_filters:
                default_jql += f' AND {" AND ".join(date_filters)}'

            jql_query = jql or f'{default_jql} order by updated desc'

            url = f'{self.jira_url}/rest/api/2/search'
            params = {
                'jql': jql_query,
                'fields': 'key,status,assignee,reporter,summary',
                'maxResults': 100
            }
            response = requests.get(url, auth=self.auth, params=params)
            response.raise_for_status()
            issues = response.json().get('issues', [])

            data = []
            for issue in issues:
                fields = issue.get('fields', {})
                data.append({
                    'issue_id': issue.get('key'),
                    'status': fields.get('status', {}).get('name', 'Unknown'),
                    'assignee': fields.get('assignee', {}).get('displayName', 'Unassigned'),
                    'reporter': fields.get('reporter', {}).get('displayName', 'Unknown'),
                    'summary': fields.get('summary', '')
                })
            return data
        except Exception as e:
            self.log.error(f"Failed to fetch JIRA issues: {e}", exc_info=True)
            return []

    def render_as_docx(self, doc: Document, data: List[Dict[str, Any]]):
        doc.add_heading("📌 JIRA Issues Worked On", level=1)

        if not data:
            doc.add_paragraph("No issues found.")
            return

        table = doc.add_table(rows=1, cols=5)
        table.style = "Table Grid"
        hdr_cells = table.rows[0].cells
        headers = ["Issue ID", "Status", "Assignee", "Reporter", "Summary"]
        for i, header in enumerate(headers):
            hdr_cells[i].text = header

        for item in data:
            row = table.add_row().cells
            row[0].text = item['issue_id']
            row[1].text = item['status']
            row[2].text = item['assignee']
            row[3].text = item['reporter']
            row[4].text = item['summary']

    def render_as_excel(self, ws: Worksheet, data: List[Dict[str, Any]]):
        headers = ["Issue ID", "Status", "Assignee", "Reporter", "Summary"]
        ws.append(headers)

        for item in data:
            ws.append([
                item['issue_id'],
                item['status'],
                item['assignee'],
                item['reporter'],
                item['summary']
            ])

        for i, _ in enumerate(headers, 1):
            ws.column_dimensions[get_column_letter(i)].width = 20

    def render_as_pdf(self, story: list, data: List[Dict[str, Any]]):
        style = getSampleStyleSheet()["Normal"]
        story.append(Paragraph("📌 JIRA Issues Worked On", getSampleStyleSheet()["Heading2"]))

        if not data:
            story.append(Paragraph("No issues found.", style))
            return

        table_data = [["Issue ID", "Status", "Assignee", "Reporter", "Summary"]]
        for item in data:
            table_data.append([
                item['issue_id'],
                item['status'],
                item['assignee'],
                item['reporter'],
                item['summary']
            ])

        table = Table(table_data, hAlign='LEFT')
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 9),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 6),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(table)

    def render_as_ppt(self, presentation: Presentation, data: List[Dict[str, Any]]):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        title_shape = slide.shapes.title
        title_shape.text = "📌 JIRA Issues Worked On"

        body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5))
        text_frame = body_shape.text_frame
        text_frame.word_wrap = True

        if not data:
            text_frame.text = "No issues found."
            return

        for item in data:
            p = text_frame.add_paragraph()
            p.text = f"{item['issue_id']} | {item['status']} | {item['assignee']} | {item['reporter']}\n{item['summary']}"
            p.level = 0

    def render_as_txt(self, data: List[Dict[str, Any]]) -> str:
        if not data:
            return "📌 JIRA Issues Worked On:\nNo issues found.\n"

        lines = ["📌 JIRA Issues Worked On:\n"]
        for item in data:
            lines.append(
                f"{item['issue_id']} | {item['status']} | {item['assignee']} | {item['reporter']}\n{item['summary']}\n")
        return "\n".join(lines)

def main(jira_url, auth_token, user_email, output_file, start_date=None, end_date=None, output_format="all"):

    logger = Logger("JiraExporterMain", logging.DEBUG)
    # Replace with your actual JIRA URL and credentials
    jira_url = jira_url
    auth = (user_email, auth_token)
    user_email = user_email

    exporter = JiraExporter(logger, jira_url, auth)

    # Fetch JIRA data
    data = exporter.generate_jira_report(user_email=user_email, start_date=start_date, end_date=end_date)
    output_file_name = f"{output_file}"

    if output_format == "docx" or output_format == "all":
        # Render to DOCX
        doc = Document()
        exporter.render_as_docx(doc, data)
        doc.save(f"{output_file_name}_jira_issues.docx")
        logger.info("Saved %s_jira_issues.docx", output_file_name)

    if output_format == "xlsx" or output_format == "all":
        # Render to Excel
        wb = Workbook()
        ws = wb.active
        exporter.render_as_excel(ws, data)
        wb.save(f"{output_file_name}_jira_issues.xlsx")
        logger.info("Saved %s_jira_issues.xlsx", output_file_name)

    if output_format == "pdf" or output_format == "all":
        # Render to PDF
        story = []
        pdf_path = f"{output_file_name}_jira_issues.pdf"
        exporter.render_as_pdf(story, data)
        story.append(Spacer(1, 12))
        SimpleDocTemplate(pdf_path, pagesize=letter).build(story)
        logger.info("Saved %s_jira_issues.pdf", output_file_name)

    if output_format == "pptx" or output_format == "all":
        # Render to PPTX
        ppt = Presentation()
        exporter.render_as_ppt(ppt, data)
        ppt.save(f"{output_file_name}_jira_issues.pptx")
        logger.info("Saved %s_jira_issues.pptx", output_file_name)

    if output_format == "txt" or output_format == "all":
        # Render to TXT
        txt_output = exporter.render_as_txt(data)
        with open(f"{output_file_name}_jira_issues.txt", "w") as f:
            f.write(txt_output)
        logger.info("Saved %s_jira_issues.txt", output_file_name)

    return 0 if data is None else len(data)