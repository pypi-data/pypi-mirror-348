# Copyright 2024 Liu Siyao
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

from __future__ import print_function

import logging
import os
from functools import partial
from operator import attrgetter
from typing import List, Union
import re

from PIL import Image
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from rapidfuzz import process as fuze_process
from tqdm import tqdm

from pptx2md.multi_column import get_multi_column_slide_if_present
from pptx2md.types import (
    ConversionConfig,
    GeneralSlide,
    ImageElement,
    ListItemElement,
    ParagraphElement,
    ParsedPresentation,
    SlideElement,
    TableElement,
    TextRun,
    TextStyle,
    TitleElement,
    VideoElement,
)

logger = logging.getLogger(__name__)

picture_count = 0
video_count = 0


def emu_to_pt(emu):
    """将EMU(English Metric Units)转换为磅值(points)"""
    return emu / 12700  # 1 pt = 12700 EMU


def is_title(shape):
    """判断一个shape是否为标题
    
    判断标准:
    1. 是标题占位符
    2. 或者是文本框且同时满足:
       - 字体大小较大(>=28pt)
       - 在幻灯片顶部(top < 100pt)
       - 文本较短(<50字符)
    """
    # 检查是否是标题占位符
    if shape.is_placeholder and (shape.placeholder_format.type == PP_PLACEHOLDER.TITLE or
                               shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE or
                               shape.placeholder_format.type == PP_PLACEHOLDER.VERTICAL_TITLE or
                               shape.placeholder_format.type == PP_PLACEHOLDER.CENTER_TITLE):
        return True
    
    # 检查是否是看起来像标题的文本框
    try:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            
            # 获取第一个段落的字体大小
            if shape.text_frame.paragraphs:
                first_para = shape.text_frame.paragraphs[0]
                if first_para.runs:
                    font_size = first_para.runs[0].font.size
                    if font_size:
                        font_size = font_size.pt
                    else:
                        font_size = 0
                else:
                    font_size = 0
            else:
                font_size = 0
            
            # 转换EMU到points
            top_pt = emu_to_pt(shape.top)
            
            # 判断条件:
            # 1. 文本不为空
            # 2. 同时满足: 字体大小>=28pt 且 在幻灯片顶部
            # 3. 文本长度适中
            if (text and  # 文本不为空
                font_size >= 28 and  # 字体大
                top_pt < 100 and  # 在顶部(100pt ≈ 3.5cm)
                len(text) < 50):  # 文本不会太长
                return True
    except Exception as e:
        logger.warning(f"Error checking title text: {e}")
    
    return False


def is_text_block(config: ConversionConfig, shape):
    """判断一个shape是否为正文文本块
    
    判断标准:
    1. 是正文占位符
    2. 或者是文本框且文本长度超过最小块大小
    """
    if not shape.has_text_frame:
        return False

    # 检查是否是正文占位符
    if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
        return True

    # 对于非占位符文本框,只检查文本长度
    text = shape.text_frame.text.strip()
    if text and len(text) > config.min_block_size:
        return True

    return False


def is_list_block(shape) -> bool:
    levels = []
    for para in shape.text_frame.paragraphs:
        if para.level not in levels:
            levels.append(para.level)
        if para.level != 0 or len(levels) > 1:
            return True
    return False


def is_accent(font):
    if font.underline or font.italic or (
            font.color.type == MSO_COLOR_TYPE.SCHEME and
        (font.color.theme_color == MSO_THEME_COLOR.ACCENT_1 or font.color.theme_color == MSO_THEME_COLOR.ACCENT_2 or
         font.color.theme_color == MSO_THEME_COLOR.ACCENT_3 or font.color.theme_color == MSO_THEME_COLOR.ACCENT_4 or
         font.color.theme_color == MSO_THEME_COLOR.ACCENT_5 or font.color.theme_color == MSO_THEME_COLOR.ACCENT_6)):
        return True
    return False


def is_strong(font):
    if font.bold or (font.color.type == MSO_COLOR_TYPE.SCHEME and (font.color.theme_color == MSO_THEME_COLOR.DARK_1 or
                                                                   font.color.theme_color == MSO_THEME_COLOR.DARK_2)):
        return True
    return False


def get_text_runs(para) -> List[TextRun]:
    runs = []
    for run in para.runs:
        result = TextRun(text=run.text, style=TextStyle())
        if result.text == '':
            continue
        try:
            if run.hyperlink.address:
                result.style.hyperlink = run.hyperlink.address
        except:
            result.style.hyperlink = 'error:ppt-link-parsing-issue'
        if is_accent(run.font):
            result.style.is_accent = True
        if is_strong(run.font):
            result.style.is_strong = True
        if run.font.color.type == MSO_COLOR_TYPE.RGB:
            rgb = run.font.color.rgb
            # 统一转为元组
            if hasattr(rgb, 'to_rgb'):
                result.style.color_rgb = rgb.to_rgb()
            else:
                # 兼容字符串形式
                hexstr = str(rgb)
                result.style.color_rgb = tuple(int(hexstr[i:i+2], 16) for i in (0, 2, 4))
        runs.append(result)
    return runs


def process_title(config: ConversionConfig, shape, slide_idx) -> TitleElement:
    text = shape.text_frame.text.strip()
    if config.custom_titles:
        res = fuze_process.extractOne(text, config.custom_titles.keys(), score_cutoff=92)
        if not res:
            return TitleElement(content=text.strip(), level=max(config.custom_titles.values()) + 1)
        else:
            logger.info(f'Title in slide {slide_idx} "{text}" is converted to "{res[0]}" as specified in title file.')
            return TitleElement(content=res[0].strip(), level=config.custom_titles[res[0]])
    else:
        return TitleElement(content=text.strip(), level=1)


def process_text_blocks(config: ConversionConfig, shape, slide_idx) -> List[Union[ListItemElement, ParagraphElement]]:
    results = []
    if is_list_block(shape):
        for para in shape.text_frame.paragraphs:
            if para.text.strip() == '':
                continue
            text = get_text_runs(para)
            results.append(ListItemElement(content=text, level=para.level))
    else:
        # paragraph block
        for para in shape.text_frame.paragraphs:
            if para.text.strip() == '':
                continue
            text = get_text_runs(para)
            results.append(ParagraphElement(content=text))
    return results


def process_picture(config: ConversionConfig, shape, slide_idx) -> Union[ImageElement, None]:
    """处理图片元素，支持普通图片和占位符中的图片"""
    if config.disable_image:
        return None

    global picture_count
    
    try:
        # 获取图片数据
        image_blob = None
        image_ext = 'png'  # 默认扩展名
        
        # 1. 直接从shape.image获取
        if hasattr(shape, 'image') and shape.image:
            image_blob = shape.image.blob
            image_ext = shape.image.ext
        
        # 2. 从占位符中获取图片
        elif hasattr(shape, 'placeholder_format') and hasattr(shape, '_element'):
            # 检查XML中是否有图片引用
            element_xml = str(shape._element)
            if 'blip' in element_xml or 'a:blip' in element_xml:
                # 尝试从关系中获取图片
                if hasattr(shape, '_parent') and hasattr(shape._parent, 'part'):
                    part = shape._parent.part
                    if hasattr(part, 'rels'):
                        for rel_id, rel in part.rels.items():
                            if hasattr(rel, 'target_part') and hasattr(rel.target_part, 'content_type'):
                                content_type = rel.target_part.content_type
                                if 'image' in content_type.lower():
                                    if hasattr(rel.target_part, 'blob'):
                                        image_blob = rel.target_part.blob
                                        # 从内容类型获取扩展名
                                        if '/' in content_type:
                                            image_ext = content_type.split('/')[-1]
                                            if image_ext == 'jpeg':
                                                image_ext = 'jpg'
                                        # 或从部件名获取扩展名
                                        elif hasattr(rel.target_part, 'partname'):
                                            part_name = str(rel.target_part.partname)
                                            if '.' in part_name:
                                                image_ext = part_name.split('.')[-1].lower()
                                        break
        
        # 如果没有找到图片数据，返回None
        if not image_blob:
            logger.warning(f"在幻灯片 {slide_idx} 中找不到图片数据")
            return None
            
        # 创建图片目录
        if not os.path.exists(config.image_dir):
            os.makedirs(config.image_dir)
            
        # 生成文件名
        file_prefix = ''.join(os.path.basename(config.pptx_path).split('.')[:-1])
        pic_name = file_prefix + f'_{slide_idx}_{picture_count}'
        output_path = config.image_dir / f'{pic_name}.{image_ext}'
        
        # 保存图片文件
        with open(output_path, 'wb') as f:
            f.write(image_blob)
        picture_count += 1
            
        # 计算相对路径
        common_path = os.path.commonpath([config.output_path, config.image_dir])
        img_outputter_path = os.path.relpath(output_path, common_path)
        
        # 处理wmf格式图片
        if image_ext == 'wmf' and not config.disable_wmf:
            try:
                try:
                    Image.open(output_path).save(os.path.splitext(output_path)[0] + '.png')
                    return ImageElement(path=os.path.splitext(img_outputter_path)[0] + '.png', width=config.image_width)
                except Exception:  # Image failed, try another
                    from wand.image import Image
                    with Image(filename=output_path) as img:
                        img.format = 'png'
                        img.save(filename=os.path.splitext(output_path)[0] + '.png')
                    logger.info(f'图片 {output_path} 在幻灯片 {slide_idx} 中转换为png格式。')
                    return ImageElement(path=os.path.splitext(img_outputter_path)[0] + '.png', width=config.image_width)
            except Exception:
                logger.warning(f'无法将wmf图片 {output_path} 在幻灯片 {slide_idx} 中转换为png格式，跳过。')
                if config.disable_wmf:
                    return None
                
        # 返回图片元素
        return ImageElement(path=img_outputter_path, width=config.image_width)
        
    except Exception as e:
        logger.warning(f'处理图片时出错，幻灯片 {slide_idx}: {e}')
        return None


def process_video(config: ConversionConfig, shape, slide_idx) -> Union[VideoElement, None]:
    """处理视频元素，提取视频数据并保存到文件"""
    if config.disable_image or config.disable_video:  # 如果禁用了图片或视频
        return None

    global video_count

    try:
        # 初始化变量
        video_ext = 'mp4'  # 默认扩展名
        video_blob = None
        blob_size = 0
        external_video_url = None
        
        # 设置日志函数
        log_fn = logger.info if config.debug else logger.debug
        log_fn(f"处理潜在视频元素，幻灯片 {slide_idx}, 形状类型: {getattr(shape, 'shape_type', 'unknown')}")
        
        # 收集所有可能的媒体部件
        media_parts = []
        
        # 1. 从shape.media属性获取
        if hasattr(shape, 'media') and shape.media:
            media = shape.media
            log_fn(f"  发现媒体属性: {media}")
            if hasattr(media, 'blob'):
                blob = media.blob
                size = len(blob) if blob else 0
                if size > 1000:  # 忽略太小的数据块
                    log_fn(f"  发现媒体数据块，大小: {size} 字节")
                    media_parts.append(('media属性', getattr(media, 'ext', 'mp4'), blob, size))
                else:
                    log_fn(f"  警告: 媒体数据块太小 ({size} 字节)，可能是缩略图或损坏的文件")
        
        # 2. 从关系中提取媒体数据
        if hasattr(shape, '_parent') and hasattr(shape._parent, 'part'):
            part = shape._parent.part
            
            # 2.1 检查直接关系
            if hasattr(part, 'rels'):
                for rel_id, rel in part.rels.items():
                    try:
                        if hasattr(rel, 'reltype') and ('media' in rel.reltype.lower() or 'video' in rel.reltype.lower()):
                            log_fn(f"  发现媒体关系: {rel_id} -> {rel.reltype}")
                            
                            # 检查是否是外部链接
                            if hasattr(rel, '_target_mode') and rel._target_mode == 'External':
                                if hasattr(rel, '_target'):
                                    external_url = rel._target
                                    log_fn(f"  发现外部视频链接: {external_url}")
                                    external_video_url = external_url
                                    continue
                            
                            # 处理嵌入式视频
                            if hasattr(rel, 'target_part') and hasattr(rel.target_part, 'blob'):
                                media_part = rel.target_part
                                part_name = str(getattr(media_part, 'partname', f'rel_{rel_id}'))
                                blob = media_part.blob
                                size = len(blob) if blob else 0
                                if size > 1000:
                                    log_fn(f"  发现媒体部件，名称: {part_name}, 大小: {size} 字节")
                                    media_parts.append((part_name, part_name.split('.')[-1].lower() if '.' in part_name else 'mp4', blob, size))
                    except Exception as e:
                        log_fn(f"  处理关系 {rel_id} 时出错: {e}")
                        # 如果错误是因为外部链接，尝试获取链接
                        if "target-mode is external" in str(e) and hasattr(rel, '_target'):
                            external_url = rel._target
                            log_fn(f"  发现外部视频链接: {external_url}")
                            external_video_url = external_url
            
            # 2.2 搜索所有包部件
            if blob_size < 5000 and hasattr(part, 'package') and hasattr(part.package, 'parts'):
                log_fn(f"  搜索所有媒体部件...")
                for pkg_part in part.package.parts:
                    part_name = str(getattr(pkg_part, 'partname', ''))
                    if any(term in part_name.lower() for term in ['media', 'video', 'movie']):
                        if hasattr(pkg_part, 'blob'):
                            blob = pkg_part.blob
                            size = len(blob) if blob else 0
                            if size > 1000:
                                log_fn(f"  发现包媒体部件: {part_name}, 大小: {size} 字节")
                                media_parts.append((part_name, part_name.split('.')[-1].lower() if '.' in part_name else 'mp4', blob, size))
        
        # 3. 从XML元素中提取关系ID
        if not media_parts and not external_video_url and hasattr(shape, '_element'):
            log_fn(f"  检查形状元素，尝试找到视频引用...")
            rel_ids = extract_media_rel_ids(shape)
            
            if rel_ids and hasattr(shape, '_parent') and hasattr(shape._parent, 'part'):
                part = shape._parent.part
                
                for rel_id in rel_ids:
                    try:
                        target_part = None
                        # 从关系获取目标部件
                        if hasattr(part, 'rels') and rel_id in part.rels:
                            rel = part.rels[rel_id]
                            
                            # 检查是否是外部链接
                            if hasattr(rel, '_target_mode') and rel._target_mode == 'External':
                                if hasattr(rel, '_target'):
                                    external_url = rel._target
                                    log_fn(f"  发现外部视频链接: {external_url}")
                                    external_video_url = external_url
                                    continue
                            
                            target_part = rel.target_part
                        # 从相关部件字典中获取
                        elif hasattr(part, 'related_parts') and rel_id in part.related_parts:
                            target_part = part.related_parts[rel_id]
                            
                        if target_part and hasattr(target_part, 'blob'):
                            part_name = str(getattr(target_part, 'partname', f'rel_{rel_id}'))
                            blob = target_part.blob
                            size = len(blob) if blob else 0
                            if size > 1000:
                                log_fn(f"  找到目标部件: {part_name}, 大小: {size} 字节")
                                media_parts.append((part_name, part_name.split('.')[-1].lower() if '.' in part_name else 'mp4', blob, size))
                    except Exception as e:
                        log_fn(f"  获取视频部件失败，关系ID {rel_id}: {e}")
                        # 如果错误是因为外部链接，尝试获取链接
                        if "target-mode is external" in str(e) and hasattr(part.rels[rel_id], '_target'):
                            external_url = part.rels[rel_id]._target
                            log_fn(f"  发现外部视频链接: {external_url}")
                            external_video_url = external_url
        
        # 4. 如果是图片形状，尝试从图片关系中查找视频
        if not media_parts and not external_video_url and hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            log_fn(f"  检查图片形状的关系，尝试找到视频...")
            if hasattr(shape, '_parent') and hasattr(shape._parent, 'part'):
                part = shape._parent.part
                
                # 列出所有关系
                if hasattr(part, 'rels'):
                    for rel_id, rel in part.rels.items():
                        try:
                            # 检查是否是外部链接
                            if hasattr(rel, '_target_mode') and rel._target_mode == 'External':
                                if hasattr(rel, '_target'):
                                    external_url = rel._target
                                    if any(ext in external_url.lower() for ext in 
                                          ['.mp4', '.mov', '.avi', '.wmv', '.mpg', '.mpeg', '.webm', '.flv']):
                                        log_fn(f"  发现外部视频链接: {external_url}")
                                        external_video_url = external_url
                                        continue
                            
                            if hasattr(rel, 'target_part'):
                                target_part = rel.target_part
                                part_name = str(getattr(target_part, 'partname', ''))
                                log_fn(f"  检查关系 {rel_id}: {part_name}")
                                
                                # 检查是否是视频文件
                                is_video_file = any(ext in part_name.lower() for ext in 
                                                  ['.mp4', '.mov', '.avi', '.wmv', '.mpg', '.mpeg', '.webm', '.flv'])
                                
                                if is_video_file and hasattr(target_part, 'blob'):
                                    blob = target_part.blob
                                    size = len(blob) if blob else 0
                                    if size > 5000:  # 确保足够大
                                        log_fn(f"  找到视频文件: {part_name}, 大小: {size} 字节")
                                        media_parts.append((part_name, part_name.split('.')[-1].lower(), blob, size))
                        except Exception as e:
                            log_fn(f"  检查关系时出错: {e}")
                            # 如果错误是因为外部链接，尝试获取链接
                            if "target-mode is external" in str(e) and hasattr(rel, '_target'):
                                external_url = rel._target
                                if any(ext in external_url.lower() for ext in 
                                      ['.mp4', '.mov', '.avi', '.wmv', '.mpg', '.mpeg', '.webm', '.flv']):
                                    log_fn(f"  发现外部视频链接: {external_url}")
                                    external_video_url = external_url
        
        # 5. 按优先级选择最佳媒体部件
        if media_parts:
            # 按优先级和大小排序
            sorted_parts = sorted(media_parts, key=lambda x: (get_media_priority(x[0], x[1]), x[3]), reverse=True)
            
            # 选择最佳部件
            best_name, best_ext, best_blob, best_size = sorted_parts[0]
            log_fn(f"  选择最佳媒体部件: {best_name}, 扩展名: {best_ext}, 大小: {best_size} 字节")
            
            video_blob = best_blob
            blob_size = best_size
            
            # 只有当扩展名是视频格式时才更新
            if best_ext in ['mp4', 'mov', 'avi', 'wmv', 'mpg', 'mpeg', 'webm', 'flv', 'm4v', 'ogg']:
                video_ext = best_ext
        
        # 6. 处理外部链接视频
        if external_video_url and not video_blob:
            log_fn(f"  使用外部视频链接: {external_video_url}")
            
            # 创建视频目录
            video_dir = config.image_dir.parent / 'video'
            os.makedirs(video_dir, exist_ok=True)
            
            # 生成文件名
            video_count += 1
            file_prefix = ''.join(os.path.basename(config.pptx_path).split('.')[:-1])
            
            # 从URL中提取扩展名
            if '.' in external_video_url:
                url_ext = external_video_url.split('.')[-1].lower()
                if url_ext in ['mp4', 'mov', 'avi', 'wmv', 'mpg', 'mpeg', 'webm', 'flv', 'm4v', 'ogg']:
                    video_ext = url_ext
            
            # 创建视频元素
            return VideoElement(
                path=external_video_url,  # 直接使用外部URL
                width=config.image_width,
                height=getattr(shape, 'height', None),
                is_external=True  # 标记为外部链接
            )
        
        # 7. 如果找到了视频数据，保存并创建视频元素
        if video_blob and blob_size > 5000:  # 确保数据块足够大
            # 创建视频目录
            video_dir = config.image_dir.parent / 'video'
            os.makedirs(video_dir, exist_ok=True)
            
            # 生成文件名
            video_count += 1
            file_prefix = ''.join(os.path.basename(config.pptx_path).split('.')[:-1])
            video_filename = f"{file_prefix}_{video_count}.{video_ext}"
            video_path = os.path.join(video_dir, video_filename)
            
            # 保存视频文件
            with open(video_path, 'wb') as f:
                f.write(video_blob)
            
            log_fn(f"  保存视频文件: {video_path}, 大小: {blob_size} 字节")
            
            # 计算相对路径
            common_path = os.path.commonpath([config.output_path, video_dir])
            video_rel_path = os.path.relpath(video_path, common_path)
            
            # 创建视频元素
            return VideoElement(
                path=video_rel_path,
                width=config.image_width,
                height=getattr(shape, 'height', None)
            )
        else:
            # 如果是图片形状但没有找到视频数据，尝试作为图片处理
            if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                log_fn(f"  未找到视频数据，尝试作为图片处理")
                return process_picture(config, shape, slide_idx)
            else:
                log_fn(f"  未找到有效的视频数据")
                return None
            
    except Exception as e:
        logger.warning(f"处理视频时出错: {e}")
        return None


def extract_media_rel_ids(shape):
    """从形状元素中提取媒体关系ID"""
    rel_ids = []
    
    try:
        element = shape._element
        from xml.etree import ElementTree
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
             'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
             'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
             'v': 'http://schemas.openxmlformats.org/drawingml/2006/video'}
        
        # 搜索videoFile标签
        video_elements = element.findall('.//p:videoFile', ns) or []
        rel_ids.extend([e.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id') 
                       for e in video_elements if e.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')])
        
        # 搜索movieFile标签
        movie_elements = element.findall('.//p:movieFile', ns) or []
        rel_ids.extend([e.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id') 
                       for e in movie_elements if e.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')])
        
        # 直接从元素中搜索关系ID
        for elem in element.iter():
            for attr in elem.attrib:
                if 'id' in attr.lower() and 'relationship' in attr.lower():
                    rel_ids.append(elem.attrib[attr])
        
        # 如果没有找到，尝试从XML文本中提取
        if not rel_ids:
            import re
            element_str = ElementTree.tostring(element, encoding='unicode')
            r_ids = re.findall(r'r:id="([^"]+)"', element_str)
            if r_ids:
                rel_ids.extend(r_ids)
            else:
                # 在整个幻灯片中搜索
                if hasattr(shape, '_parent') and hasattr(shape._parent, 'element'):
                    slide_element = shape._parent.element
                    slide_str = ElementTree.tostring(slide_element, encoding='unicode')
                    slide_r_ids = re.findall(r'r:id="([^"]+)"', slide_str)
                    if slide_r_ids:
                        for rid in slide_r_ids:
                            if 'video' in rid.lower() or 'media' in rid.lower():
                                rel_ids.append(rid)
    except Exception as e:
        logger.warning(f"提取媒体关系ID时出错: {e}")
    
    return rel_ids


def get_media_priority(part_name, ext):
    """定义媒体文件的优先级"""
    # 按文件扩展名优先级
    ext_priority = {
        'mp4': 10,
        'mov': 9,
        'avi': 8,
        'wmv': 7,
        'mpg': 6, 'mpeg': 6,
        'webm': 5,
        'flv': 4,
        'm4v': 3,
        'ogg': 2
    }
    
    # 获取扩展名优先级
    priority = ext_priority.get(ext.lower(), 0)
    
    # 根据部件名称增加优先级
    if 'video' in part_name.lower():
        priority += 2
    elif 'media' in part_name.lower():
        priority += 1
    elif 'movie' in part_name.lower():
        priority += 1
        
    # 排除XML文件
    if ext.lower() == 'xml':
        priority = -1
        
    return priority


def process_table(config: ConversionConfig, shape, slide_idx) -> Union[TableElement, None]:
    """处理表格元素"""
    try:
        table = shape.table
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                text_runs = []
                for paragraph in cell.text_frame.paragraphs:
                    text_runs.extend(get_text_runs(paragraph))
                cells.append(text_runs)  # 每个单元格是TextRun列表
            rows.append(cells)
        return TableElement(content=rows)
    except Exception as e:
        logger.warning(f"处理表格时出错，幻灯片 {slide_idx}: {e}")
    return None


def ungroup_shapes(shapes) -> List[SlideElement]:
    res = []
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                res.extend(ungroup_shapes(shape.shapes))
            else:
                res.append(shape)
        except Exception as e:
            logger.warning(f'failed to load shape {shape}, skipped. error: {e}')
    return res


def process_shapes(config: ConversionConfig, current_shapes, slide_id: int) -> List[SlideElement]:
    """处理幻灯片中的所有形状，提取文本、图片和表格"""
    results = []
    
    logger.info(f"处理幻灯片 {slide_id} 中的 {len(current_shapes)} 个形状")
    
    for i, shape in enumerate(current_shapes):
        try:
            shape_type = getattr(shape, 'shape_type', 'unknown')
            shape_name = getattr(shape, 'name', 'unnamed')
            logger.info(f"  处理形状 {i+1}/{len(current_shapes)}: {shape_name} (类型: {shape_type})")
            
            # 处理标题
            if is_title(shape):
                logger.info(f"    识别为标题")
                title = process_title(config, shape, slide_id)
                if title:
                    results.append(title)
                    
            # 处理文本块
            elif is_text_block(config, shape):
                logger.info(f"    识别为文本块")
                text_elements = process_text_blocks(config, shape, slide_id)
                if text_elements:
                    results.extend(text_elements)
                    
            # 处理表格
            elif hasattr(shape, 'table'):
                logger.info(f"    识别为表格")
                table = process_table(config, shape, slide_id)
                if table:
                    results.append(table)
                    
            # 处理组合形状
            elif hasattr(shape, 'shapes'):
                logger.info(f"    识别为组合形状，包含 {len(shape.shapes)} 个子形状")
                sub_elements = process_shapes(config, shape.shapes, slide_id)
                if sub_elements:
                    results.extend(sub_elements)
                    
            # 处理视频
            elif is_video_shape(shape):
                logger.info(f"    识别为视频")
                try:
                    video = process_video(config, shape, slide_id)
                    if video:
                        logger.info(f"    成功提取视频: {video.path}" + (" (外部链接)" if getattr(video, 'is_external', False) else ""))
                        results.append(video)
                    else:
                        logger.warning(f"    视频提取失败，尝试作为图片处理")
                        try_as_picture = True
                except Exception as e:
                    logger.warning(f'    处理嵌入视频失败: {e}')
                    # 尝试作为图片处理
                    try_as_picture = True
                
                # 如果视频提取失败，尝试作为图片处理
                if 'try_as_picture' in locals() and try_as_picture:
                    try:
                        # 检查是否有图片特征
                        if is_picture_shape(shape) or hasattr(shape, 'image') or shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                            logger.info(f"    尝试作为图片处理")
                            picture = process_picture(config, shape, slide_id)
                            if picture:
                                logger.info(f"    成功提取图片: {picture.path}")
                                results.append(picture)
                        else:
                            logger.warning(f"    图片提取也失败")
                    except Exception as pic_e:
                        logger.warning(f'    处理图片也失败: {pic_e}')
                    
            # 处理图片
            elif is_picture_shape(shape):
                logger.info(f"    识别为图片")
                picture = process_picture(config, shape, slide_id)
                if picture:
                    results.append(picture)
                    
            # 处理占位符
            elif shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                logger.info(f"    识别为占位符")
                # 检查占位符类型
                try:
                    if hasattr(shape, 'placeholder_format'):
                        ph_type = shape.placeholder_format.type
                        logger.info(f"    占位符类型: {ph_type}")
                        
                        # 处理媒体占位符
                        if ph_type in [PP_PLACEHOLDER.MEDIA_CLIP, PP_PLACEHOLDER.MEDIA]:
                            # 首先尝试作为视频处理
                            try:
                                logger.info(f"    尝试作为视频处理")
                                video = process_video(config, shape, slide_id)
                                if video:
                                    logger.info(f"    成功提取视频: {video.path}" + (" (外部链接)" if getattr(video, 'is_external', False) else ""))
                                    results.append(video)
                                    continue
                            except Exception as e:
                                logger.warning(f"    视频处理失败: {e}")
                            
                            # 如果视频处理失败，尝试作为图片处理
                            logger.info(f"    尝试作为图片处理")
                            try:
                                picture = process_picture(config, shape, slide_id)
                                if picture:
                                    logger.info(f"    成功提取图片: {picture.path}")
                                    results.append(picture)
                                    continue
                            except Exception as pic_e:
                                logger.warning(f"    图片处理失败: {pic_e}")
                except Exception as e:
                    logger.warning(f"    处理占位符时出错: {e}")
                    # 如果错误信息中包含"MEDIA"，尝试作为图片处理
                    if "MEDIA" in str(e):
                        logger.info(f"    占位符错误中包含MEDIA关键字，尝试作为图片处理")
                        try:
                            picture = process_picture(config, shape, slide_id)
                            if picture:
                                logger.info(f"    成功提取图片: {picture.path}")
                                results.append(picture)
                        except Exception as pic_e:
                            logger.warning(f"    图片处理失败: {pic_e}")
                
        except Exception as e:
            logger.warning(f'    处理形状失败: {e}')

    logger.info(f"幻灯片 {slide_id} 处理完成，提取了 {len(results)} 个元素")
    return results


def is_picture_shape(shape) -> bool:
    """检测一个形状是否是图片"""
    if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    return False


def is_video_shape(shape) -> bool:
    """更准确地检测一个形状是否包含视频"""
    try:
        shape_id = getattr(shape, 'id', 'unknown')
        shape_name = getattr(shape, 'name', 'unnamed')
        logger.info(f"检查形状是否为视频: ID={shape_id}, 名称={shape_name}")
        
        # 1. 快速排除明确的非视频类型
        if hasattr(shape, 'shape_type'):
            shape_type = shape.shape_type
            if shape_type in [MSO_SHAPE_TYPE.TABLE, MSO_SHAPE_TYPE.GROUP]:
                logger.info(f"  [不是视频] 形状类型是 {shape_type}")
                return False
            
            # 对于MEDIA类型，直接认为是视频
            if shape_type == MSO_SHAPE_TYPE.MEDIA:
                logger.info(f"  [是视频] 形状类型是 MEDIA")
                return True
            
            # 对于图片类型，需要更严格的检查
            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                # 检查图片名称是否包含视频相关词汇
                if hasattr(shape, 'name') and any(term in shape.name.lower() for term in ['video', 'movie', '视频', '影片', '对话']):
                    logger.info(f"  [可能是视频] 图片名称包含视频相关词汇: {shape.name}")
                else:
                    # 检查是否有明确的视频特征
                    has_video_feature = False
                    
                    # 检查是否有media属性
                    if hasattr(shape, 'media') and shape.media is not None:
                        has_video_feature = True
                        logger.info(f"  [是视频] 图片形状有media属性")
                    
                    # 检查XML中是否有视频标记
                    elif hasattr(shape, '_element'):
                        element_xml = str(shape._element)
                        video_indicators = ['videoFile', 'movieFile', 'embeddedVideo', 'p:video', 'p:movie']
                        for indicator in video_indicators:
                            if indicator in element_xml:
                                has_video_feature = True
                                logger.info(f"  [是视频] 图片XML中找到视频指示器: {indicator}")
                                break
                    
                    # 如果没有明确的视频特征，则认为是普通图片
                    if not has_video_feature:
                        logger.info(f"  [不是视频] 形状是普通图片，没有视频特征")
                        return False
                
        # 2. 排除有表格或图表属性的形状
        if hasattr(shape, 'table') or hasattr(shape, 'chart'):
            logger.info(f"  [不是视频] 形状有table或chart属性")
            return False
        
        # 3. 快速确认明确的视频类型
        if hasattr(shape, 'media') and shape.media is not None:
            logger.info(f"  [是视频] 形状有media属性")
            return True
            
        # 4. 检查XML元素中的表格和视频标记
        if hasattr(shape, '_element'):
            element = shape._element
            element_xml = str(element)
            logger.info(f"  检查XML元素: 标签={element.tag if hasattr(element, 'tag') else 'unknown'}")
            
            # 排除表格相关标记
            table_indicators = ['<a:tbl', '<a:tr', '<a:tc', '<a:gridCol']
            for indicator in table_indicators:
                if indicator in element_xml:
                    logger.info(f"  [不是视频] XML中包含表格标记: {indicator}")
                    return False
            
            # 检查视频相关标记
            video_indicators = [
                'videoFile', 'movieFile', 'embeddedVideo', 'nvPr',
                'p:video', 'p:movie', 'a:videoFile', 'r:video', 'p:nvPr',
                'media/media', '.mp4', '.avi', '.mov', '.wmv',
                'video.bin', 'media.bin', 'media/video', 'embeddings/video',
                'http://schemas.openxmlformats.org/drawingml/2006/video',
                'http://schemas.microsoft.com/office/2007/relationships/media'
            ]
            
            for indicator in video_indicators:
                if indicator in element_xml:
                    logger.info(f"  [是视频] XML中找到视频指示器: {indicator}")
                    return True
            
            logger.info(f"  XML中未发现视频指示器")
                    
        # 5. 检查占位符类型
        try:
            if hasattr(shape, 'placeholder_format'):
                ph_type = shape.placeholder_format.type
                logger.info(f"  检查占位符类型: {ph_type}")
                
                # 排除表格和图表占位符
                if ph_type in [PP_PLACEHOLDER.TABLE, PP_PLACEHOLDER.CHART]:
                    logger.info(f"  [不是视频] 占位符类型是表格或图表")
                    return False
                
                # 确认媒体占位符
                if ph_type in [PP_PLACEHOLDER.MEDIA_CLIP, PP_PLACEHOLDER.MEDIA]:
                    logger.info(f"  [是视频] 占位符类型表明这是媒体元素")
                    return True
        except Exception as e:
            # 如果错误信息中包含"MEDIA"，可能是视频
            if "MEDIA" in str(e) and not (hasattr(shape, 'table') or 
                                         (hasattr(shape, '_element') and '<a:tbl' in str(shape._element))):
                logger.info(f"  [是视频] 占位符类型错误中包含MEDIA关键字: {e}")
                return True
            logger.info(f"  检查占位符类型时出错: {e}")
                
        # 6. 检查形状名称
        if hasattr(shape, 'name') and shape.name:
            name = shape.name.lower()
            
            # 排除表格相关名称
            for term in ['table', '表格', 'tbl', 'grid']:
                if term in name:
                    logger.info(f"  [不是视频] 形状名称包含表格相关词汇: {term}")
                    return False
            
            # 确认媒体相关名称
            for term in ['video', 'movie', 'media', '视频', '影片', '媒体', '对话']:
                if term in name:
                    logger.info(f"  [是视频] 形状名称包含媒体相关词汇: {term}")
                    return True
                
        # 7. 检查元素的尺寸比例 - 对于PICTURE类型，仅当有其他视频特征时才考虑比例
        if hasattr(shape, 'width') and hasattr(shape, 'height') and shape.width > 0 and shape.height > 0:
            try:
                ratio = shape.width / shape.height
                logger.info(f"  检查形状比例: {ratio:.2f} (宽={shape.width}, 高={shape.height})")
                
                # 排除表格特征
                if hasattr(shape, 'table') or (hasattr(shape, '_element') and '<a:tbl' in str(shape._element)):
                    logger.info(f"  [不是视频] 形状具有表格特征")
                    return False
                
                # 对于图片类型，仅当有明确的视频特征时才考虑比例
                if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # 检查是否有明确的视频特征
                    has_video_feature = False
                    if hasattr(shape, 'media') and shape.media is not None:
                        has_video_feature = True
                    elif hasattr(shape, '_element') and any(indicator in str(shape._element) for indicator in 
                                                          ['videoFile', 'movieFile', 'p:video', 'p:movie']):
                        has_video_feature = True
                    elif hasattr(shape, 'name') and any(term in shape.name.lower() for term in 
                                                      ['video', 'movie', 'media', '视频', '影片', '媒体', '对话']):
                        has_video_feature = True
                    
                    if not has_video_feature:
                        logger.info(f"  [不是视频] 虽然比例可能符合视频格式，但图片没有明确的视频特征")
                        return False
                
                # 检查是否符合视频比例 - 放宽比例限制
                is_video_ratio = (0.8 <= ratio < 2.5) or (0.3 < ratio <= 1.2)
                if is_video_ratio:
                    if ratio > 0.8:
                        logger.info(f"  形状比例符合横屏视频格式 (0.8 <= {ratio:.2f} < 2.5)")
                    else:
                        logger.info(f"  形状比例符合竖屏或正方形视频格式 (0.3 < {ratio:.2f} <= 1.2)")
                    
                    # 检查关系中是否有媒体引用
                    if check_video_relationships(shape, logger.info):
                        return True
                    
                    # 如果形状名称包含"对话"或其他视频相关词汇，即使没有找到媒体关系，也认为是视频
                    if hasattr(shape, 'name') and any(term in shape.name.lower() for term in 
                                                    ['video', 'movie', 'media', '视频', '影片', '媒体', '对话']):
                        logger.info(f"  [是视频] 形状名称包含媒体相关词汇且比例符合视频格式")
                        return True
            except Exception as e:
                logger.info(f"  检查形状比例时出错: {e}")
                
        # 8. 最后检查：如果是MEDIA类型或名称包含"对话"，即使前面的检查都不符合，也认为是视频
        if (hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.MEDIA) or \
           (hasattr(shape, 'name') and '对话' in shape.name.lower()):
            logger.info(f"  [是视频] 形状是MEDIA类型或名称包含'对话'")
            return True
                
        logger.info(f"  [不是视频] 没有足够的视频特征")
        return False
    except Exception as e:
        logger.warning(f"检测视频形状时出错: {e}")
        return False


def check_video_relationships(shape, log_fn):
    """检查形状的关系中是否有视频引用"""
    # 可能是视频，检查其他特征
    if hasattr(shape, '_parent') and hasattr(shape._parent, 'part') and hasattr(shape._parent.part, 'rels'):
        log_fn(f"  检查幻灯片关系...")
        found_media_rels = False
        for rel in shape._parent.part.rels:
            target = shape._parent.part.rels[rel].target_ref
            if isinstance(target, str):
                log_fn(f"    关系 {rel}: {target}")
                if 'media' in target.lower() or 'video' in target.lower():
                    log_fn(f"  [是视频] 在幻灯片关系中找到媒体引用: {target}")
                    found_media_rels = True
        if found_media_rels:
            return True
    else:
        log_fn(f"  形状没有父级部件或关系属性")
    return False


def parse(config: ConversionConfig, prs: Presentation) -> ParsedPresentation:
    """解析PowerPoint文件，提取文本、图片和表格"""
    global picture_count, video_count
    picture_count = 0
    video_count = 0
    
    # 设置日志函数
    log_fn = logger.info if config.debug else logger.debug
    
    # 创建图片目录
    if not config.disable_image and config.image_dir:
        os.makedirs(config.image_dir, exist_ok=True)
    
    # 读取自定义标题
    if config.title_path:
        with open(config.title_path, 'r', encoding='utf-8') as f:
            for i, line in enumerate(f):
                config.custom_titles[line.strip()] = i + 1
    
    # 初始化幻灯片标题字典
    config.slide_titles = {}
    
    # 解析幻灯片
    slides = []
    for i, slide in enumerate(prs.slides):
        # 如果指定了页码，只处理该页
        if config.page is not None and i + 1 != config.page:
            continue
            
        log_fn(f"========== 处理幻灯片 {i+1} ==========")
        
        # 提取幻灯片标题
        slide_title = ""
        for shape in slide.shapes:
            if is_title(shape):
                try:
                    slide_title = shape.text
                    break
                except Exception:
                    pass
        
        log_fn(f"幻灯片标题: {slide_title}")
        config.slide_titles[i+1] = slide_title
        
        # 处理幻灯片中的形状
        elements = process_shapes(config, slide.shapes, i+1)
        
        # 提取演讲者备注
        notes = []
        if not config.disable_notes and hasattr(slide, 'notes_slide') and slide.notes_slide:
            for note_shape in slide.notes_slide.notes_text_frame.paragraphs:
                note_text = note_shape.text.strip()
                if note_text:
                    notes.append(note_text)
        
        # 创建幻灯片对象
        slide_obj = GeneralSlide(elements=elements, notes=notes)
        slides.append(slide_obj)
    
    # 创建演示文稿对象
    presentation = ParsedPresentation(slides=slides)
    
    return presentation
